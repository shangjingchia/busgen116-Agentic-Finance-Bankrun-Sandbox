"""
Export the demo deck to PowerPoint (SLIDES.pptx).

Generates the demo deck — 8 main slides + B1–B8 appendix, dark palette — as an editable .pptx.
This script is the source of truth for SLIDES.pptx; re-running it overwrites the file,
so edit the text here (not just the .pptx) to keep changes durable.
One-off tooling; python-pptx is NOT a project dependency.

Run:  .venv\\Scripts\\python.exe scripts/export_slides_pptx.py
Out:  SLIDES.pptx in the project root.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x18)
DEMO_BG = RGBColor(0x24, 0x1C, 0x2E)
INK     = RGBColor(0xFF, 0xFF, 0xFF)
BODY    = RGBColor(0xD7, 0xD8, 0xE6)
RED     = RGBColor(0xE1, 0x57, 0x59)
TEAL    = RGBColor(0x76, 0xB7, 0xB2)
MUTED   = RGBColor(0x9A, 0xA0, 0xB5)
DIM     = RGBColor(0x5A, 0x5B, 0x70)
FONT    = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _box(slide, left, top, width, height, anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _para(tf, segments, size, *, bold=False, italic=False, align=PP_ALIGN.LEFT,
          space_after=8, line=None, first=False):
    """Add a paragraph. `segments` is a str or a list of (text, color) runs."""
    if first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    if isinstance(segments, str):
        segments = [(segments, INK)]
    for text, color in segments:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
        r.font.color.rgb = color
    return p


def kicker(slide, text):
    tf = _box(slide, 0.9, 0.6, 11.5, 0.5)
    _para(tf, [(text.upper(), RED)], 14, bold=True, first=True)


def footer(slide, left_text, page):
    tf = _box(slide, 0.9, 6.95, 9.0, 0.4)
    _para(tf, [(left_text, DIM)], 12, first=True)
    tf2 = _box(slide, 11.4, 6.95, 1.0, 0.4)
    _para(tf2, [(str(page), DIM)], 12, align=PP_ALIGN.RIGHT, first=True)


def bullets(tf, items, size=20, color=BODY, space_after=12):
    for i, it in enumerate(items):
        segs = it if isinstance(it, list) else [("•  ", RED), (it, color)]
        _para(tf, segs, size, space_after=space_after, line=1.25, first=(i == 0))


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = add_slide()
kicker(s, "Agentic Finance · Stress Sandbox")
tf = _box(s, 0.9, 1.8, 11.5, 2.6)
_para(tf, [("When your money runs at\n", INK)], 54, bold=True, first=True, space_after=0)
_para(tf, [("machine speed", RED)], 54, bold=True, space_after=0)
tf2 = _box(s, 0.9, 4.6, 11.5, 1.2)
_para(tf2, [("A sandbox for AI-delegated bank runs — 12 LLM agents, real decisions, one rumor.", MUTED)],
      22, first=True)
footer(s, "AI Bank Run Sandbox", "Live demo follows")

# ── Slide 2 — The shift ──────────────────────────────────────────────────────
s = add_slide()
tf = _box(s, 0.9, 0.9, 11.5, 1.4)
_para(tf, [("In 5 years, ", INK), ("an agent", RED), (" moves your money — not you.", INK)],
      36, bold=True, first=True, line=1.1)
tf2 = _box(s, 0.9, 2.7, 11.5, 4.0)
_para(tf2, [("A human bank run is slow because of ", INK), ("attention and hesitation", TEAL),
            (" — you notice, second-guess, decide. Most people aren't watching "
             "at 3am.", INK)], 22, first=True, space_after=16, line=1.35)
_para(tf2, [("Even SVB, which ran entirely online, took most of "
             "a day — that hesitation is part of what stops a rumor from becoming a collapse.", BODY)], 22,
      space_after=16, line=1.35)
_para(tf2, [("AI delegates remove it — they watch every feed, act in seconds, "
             "day or night, and tend to reason alike.", BODY)], 22, space_after=22, line=1.35)
_para(tf2, [("What does that do to a run?", RED)], 26, bold=True)
footer(s, "The premise", 2)

# ── Slide 3 — The setup ──────────────────────────────────────────────────────
s = add_slide()
tf = _box(s, 0.9, 0.8, 11.5, 1.0)
_para(tf, [("The sandbox", INK)], 36, bold=True, first=True)
tf2 = _box(s, 0.9, 1.8, 11.5, 0.9)
_para(tf2, [("Not a forecast. Not an equation. 12 AI agents making ", INK),
            ("real LLM calls", TEAL), (" with simulated money.", INK)], 22, first=True, line=1.3)
tf3 = _box(s, 0.9, 2.8, 11.5, 0.8)
_para(tf3, [("🧓 Cautious retiree     📈 Aggressive trader     🚗 Gig worker     "
             "🏛️ Institutional treasurer", BODY)], 18, first=True)
tf4 = _box(s, 0.9, 3.7, 11.5, 3.0)
bullets(tf4, [
    "Two banks. A rumor hits Bank A — with tunable credibility.",
    "Each agent sees the rumor, sees what peers do, and reasons with stakes.",
    [("•  It decides: ", RED), ("hold · partially withdraw · run.", TEAL)],
    "We watch the reasoning — verbatim — and the cascade unfold.",
], size=20)
footer(s, "How it works", 3)

# ── Slide 4 — LIVE DEMO marker ───────────────────────────────────────────────
s = add_slide(DEMO_BG)
tf = _box(s, 0.9, 2.0, 11.5, 0.6)
_para(tf, [("SWITCH TO DASHBOARD", RED)], 16, bold=True, align=PP_ALIGN.CENTER, first=True)
tf2 = _box(s, 0.9, 2.7, 11.5, 1.8)
_para(tf2, [("▶ LIVE DEMO", INK)], 72, bold=True, align=PP_ALIGN.CENTER, first=True)
tf3 = _box(s, 0.9, 4.7, 11.5, 0.8)
_para(tf3, [("Run a scenario · read an AI's mind · walk the findings", MUTED)], 22,
      align=PP_ALIGN.CENTER, first=True)
footer(s, "localhost:8501", 4)

# ── Slide 5 — SVB anchor ─────────────────────────────────────────────────────
s = add_slide()
kicker(s, "For scale — not a model of it")
tf = _box(s, 0.9, 1.2, 11.5, 1.0)
_para(tf, [("Silicon Valley Bank, March 2023", INK)], 34, bold=True, first=True)
tf2 = _box(s, 0.9, 2.1, 11.5, 1.6)
_para(tf2, [("$42B", RED)], 88, bold=True, first=True)
tf3 = _box(s, 0.9, 3.9, 11.5, 2.6)
_para(tf3, [("withdrawn in a single day — the fastest run in modern history.", INK)], 24,
      first=True, space_after=16, line=1.3)
_para(tf3, [("Why so fast? Depositors were all startups, advised by the same VCs, watching the "
             "same threads — ", BODY), ("correlated, networked, instant.", TEAL)], 22,
      space_after=16, line=1.3)
_para(tf3, [("AI delegation manufactures that correlation — for any bank — and adds a new axis: "
             "the same model.", RED)], 24, bold=True)
footer(s, "Empirical anchor", 5)

# ── Slide 6 — The numbers (model headline + speed) ───────────────────────────
s = add_slide()
kicker(s, "Two numbers to walk out with")
tf = _box(s, 0.9, 1.1, 11.5, 1.2)
_para(tf, [("Same healthy bank, same rumor — change only the ", INK), ("AI model", RED), (":", INK)],
      28, bold=True, first=True, line=1.15)
tf2 = _box(s, 0.9, 2.3, 11.5, 1.4)
_para(tf2, [("50% → 100%", RED)], 76, bold=True, first=True)
tf3 = _box(s, 0.9, 3.8, 11.5, 0.9)
_para(tf3, [("of depositors run, decided by ", INK), ("nothing but the model", TEAL),
            (" — replicated, bands don't overlap. A monoculture panics in lockstep: the model is a "
             "systemic risk factor.", INK)], 18, first=True, line=1.3)
tf4 = _box(s, 0.9, 5.0, 11.5, 1.7)
_para(tf4, [("① The model swings the outcome — ", RED), ("Claude ~50% · GPT 100% on an "
            "identical healthy bank; none tells a real crisis from a fake one well.", BODY)], 18,
      first=True, space_after=10, line=1.3)
_para(tf4, [("② …and it's ≈3× faster — ", RED), ("half decide to pull out in ~5s "
            "vs ~15s; the bank freezes first, 5/12 vs 11/12 get cash. Faster, not bigger.", BODY)], 18,
      line=1.3)
footer(s, "The headline", 6)

# ── Slide 7 — Takeaway ───────────────────────────────────────────────────────
s = add_slide()
tf = _box(s, 0.9, 0.8, 11.5, 1.4)
_para(tf, [("What we'd tell anyone ", INK), ("building or regulating", RED),
           (" AI money-managers", INK)], 32, bold=True, first=True, line=1.1)
tf2 = _box(s, 0.9, 2.4, 11.5, 3.4)
bullets(tf2, [
    [("•  ", RED), ("The model is a systemic risk factor", INK),
     (" — one model swung a healthy bank from a 50% scare to a 100% run; a fleet on one model "
      "panics in lockstep.", BODY)],
    [("•  ", RED), ("Speed compresses the response window to seconds", INK),
     (" — human oversight can't keep up.", BODY)],
    [("•  ", RED), ("Agents react to alarming language, not stated credibility", INK),
     (" — and their confidence is just as high when they're wrong.", BODY)],
], size=21, space_after=18)
tf3 = _box(s, 0.9, 5.9, 11.5, 0.9)
_para(tf3, [("A 12-agent sandbox — built to find these failure modes cheaply, before this is "
             "how finance works.", MUTED)], 18, italic=True, first=True, line=1.3)
footer(s, "One more thing →", 7)

# ── Slide 8 — Sandbox (closing recording) ────────────────────────────────────
s = add_slide(DEMO_BG)
kicker(s, "Play recording · the Sandbox")
tf = _box(s, 0.9, 1.2, 11.5, 1.1)
_para(tf, [("It isn't a fixed demo — ", INK), ("build any crisis", RED)], 34, bold=True, first=True)
# Video placeholder — replace the capability bullets (now spoken in the script) with a
# 16:9 target you drop the Sandbox recording onto in PowerPoint, then delete this box.
_VW, _VH = 6.0, 3.375  # 16:9
ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches((13.333 - _VW) / 2), Inches(2.0), Inches(_VW), Inches(_VH))
ph.fill.solid()
ph.fill.fore_color.rgb = RGBColor(0x17, 0x11, 0x1E)
ph.line.color.rgb = TEAL
ph.line.width = Pt(1.25)
ph.shadow.inherit = False
phtf = ph.text_frame
phtf.word_wrap = True
phtf.vertical_anchor = MSO_ANCHOR.MIDDLE
_para(phtf, [("▶  Sandbox screen recording", MUTED)], 20, bold=True,
      align=PP_ALIGN.CENTER, first=True, space_after=4)
_para(phtf, [("insert video here · Playback ▸ Play Full Screen · then delete this box", DIM)], 12,
      align=PP_ALIGN.CENTER)
tf3 = _box(s, 0.9, 5.6, 11.5, 1.2)
_para(tf3, [("It's not just a demo — it's ", INK), ("infrastructure", TEAL),
            (": a regulator stress-tests delegation rules before approving them, a fintech vets "
             "its agent before launch, and real data on who delegates to which model sizes the "
             "systemic risk.", INK)], 17, first=True, line=1.22)
footer(s, "Thank you — questions?", 8)

# =============================================================================
#  BACKUP / APPENDIX SLIDES — flip to these to answer Q&A. Not shown in the talk.
# =============================================================================

BACKUP_BG = RGBColor(0x15, 0x15, 0x22)


def backup_slide(title_segs, tag):
    s = add_slide(BACKUP_BG)
    tf = _box(s, 0.9, 0.55, 11.5, 0.5)
    _para(tf, [("BACKUP · FOR Q&A", MUTED)], 13, bold=True, first=True)
    tt = _box(s, 0.9, 1.12, 11.5, 1.1)
    segs = title_segs if isinstance(title_segs, list) else [(title_segs, INK)]
    _para(tt, segs, 30, bold=True, first=True, line=1.1)
    tf2 = _box(s, 11.0, 6.95, 1.4, 0.4)
    _para(tf2, [(tag, DIM)], 12, align=PP_ALIGN.RIGHT, first=True)
    return s


# B0 — divider
s = add_slide(BACKUP_BG)
tf = _box(s, 0.9, 3.0, 11.5, 1.4, anchor=MSO_ANCHOR.MIDDLE)
_para(tf, [("Appendix", INK)], 48, bold=True, first=True)
_para(tf, [("Backup slides for Q&A — flip here as needed.", MUTED)], 22, space_after=0)

# B1 — Model finding, the detail (the headline — replicated + discrimination)
s = backup_slide([("Same personas, different brains — ", INK), ("the detail", RED)], "B1")
tf = _box(s, 0.9, 2.3, 11.5, 3.0)
bullets(tf, [
    [("•  Healthy bank, withdrawing = wrong. Ran on it (mean of 4–5 reps): ", BODY)],
    [("     Claude 50% · Grok 67% · DeepSeek 79% · Mistral 85% · Gemini 98% · ", INK), ("GPT 100%", RED)],
    [("•  Bands ", BODY), ("don't overlap", INK), (" — a 50-pt gap that survives replication, not sampling noise.", BODY)],
    [("•  Discrimination (real vs. fake crisis): ", BODY), ("GPT/Gemini ≈0", INK), (" (cry-wolf, run either way); "
     "Claude +20pts is the only meaningful discriminator — and still panics 50% on a false alarm.", BODY)],
    [("•  Excluded honestly: ", BODY), ("Qwen, ByteDance-Seed", INK), (" can't emit a structured decision "
     "(default to hold) — a 0% there is plumbing, not calm.", BODY)],
], size=18, space_after=11)
tf2 = _box(s, 0.9, 5.5, 11.5, 1.1)
_para(tf2, [("Same persona can privately reach “79% chance of trouble” under both models — GPT acts on it, "
            "Claude flags its own overconfidence and hedges. That self-doubt is the whole gap.", INK)],
      17, first=True, italic=True, line=1.3)

# B2 — Speed, precisely (money-movement + online-banking answer)
s = backup_slide([("Speed, precisely — ", INK), ("decision vs. settlement", RED)], "B2")
tf = _box(s, 0.9, 2.3, 11.5, 3.4)
bullets(tf, [
    [("•  Half the depositors ", BODY), ("decide & submit", INK), (":  AI ~5s · human ~15s  → ~3×", BODY)],
    [("•  First withdrawal:  ", BODY), ("AI ~2.5s · human ~5.5s", INK)],
    [("•  Bank ", BODY), ("freezes", INK), (":  AI ~5s · human ~43s", BODY)],
    [("•  Half the cash actually ", BODY), ("settles", INK), (":  never — the bank froze first (both)", BODY)],
    [("•  Got paid before the freeze:  ", BODY), ("AI 5/12 · human 11/12", INK)],
], size=20, space_after=12)
tf2 = _box(s, 0.9, 5.7, 11.5, 1.1)
_para(tf2, [("AI speeds the ", INK), ("decision", RED), (", not the banking rail. Settlement is the "
            "same rate-limited rail for everyone — which is exactly why the faster request wave "
            "freezes the bank sooner and locks more people out.", INK)], 18, first=True, line=1.3)

# B2 — How the agents decide (not scripted)
s = backup_slide([("How the agents decide — ", INK), ("not scripted", RED)], "B3")
tf = _box(s, 0.9, 2.3, 11.5, 4.2)
bullets(tf, [
    [("•  Every decision is a ", BODY), ("real LLM call", INK), (" — persona + portfolio + what they "
     "observe in, structured action out. No utility function, no lookup table.", BODY)],
    [("•  Each persona carries a ", BODY), ("qualitative cost function", INK), (": “losing principal "
     "is catastrophic; an early-withdrawal fee is minor.” Never a numeric “withdraw if X” threshold.", BODY)],
    [("•  We deliberately ", BODY), ("don't tell them when to run", INK), (" — the model weighs the judgment call.", BODY)],
    [("•  Proof it isn't scripted: two agents of the ", BODY), ("same archetype", INK), (", same data, "
     "reach opposite calls — visible in the Inspect view.", BODY)],
], size=20, space_after=14)

# B3 — What else we found (the findings we cut from the page)
s = backup_slide([("What else we found ", INK), ("(beyond the headline three)", MUTED)], "B4")
tf = _box(s, 0.9, 2.3, 11.5, 4.2)
bullets(tf, [
    [("•  ", RED), ("Credibility label ignored", INK), (": withdrawals don't track the stated % — and a "
     "rumor labelled just 5% credible still cascaded the bank (9/12 ran).", BODY)],
    [("•  ", RED), ("Homogeneity amplifies", INK), (": all-retiree and all-treasurer populations cascade "
     "harder than a mixed one — disagreement absorbs the panic.", BODY)],
    [("•  ", RED), ("Overconfident when wrong", INK), (": agents were just as sure fleeing a healthy bank "
     "as when they were right — confidence is not a safety signal.", BODY)],
    [("•  ", RED), ("Walk-back differs by model", INK), (": GPT locks its panic in (0% reversal); Claude "
     "reconsiders and reverses ~50% — same doubt, opposite follow-through.", BODY)],
], size=19, space_after=13)

# B4 — How it's built
s = backup_slide("How it's built", "B5")
tf = _box(s, 0.9, 2.3, 11.5, 4.2)
bullets(tf, [
    [("•  ", RED), ("Event-driven engine", INK), (" (heapq queue), ~12 agents, async LLM calls fanned "
     "out so a full run is seconds, not minutes.", BODY)],
    [("•  ", RED), ("OpenRouter", INK), (" (OpenAI-compatible) — lets us swap the underlying model; Haiku "
     "for routine decisions, Sonnet for strategic moments.", BODY)],
    [("•  ", RED), ("Every decision is a structured tool call, fully logged", INK), (" — that log is the "
     "Inspect view and the audit trail.", BODY)],
    [("•  ", RED), ("Runs saved as JSON", INK), (", replayable; cost is ", BODY), ("under $1 per run.", INK)],
], size=20, space_after=14)

# B5 — Empirical anchors
s = backup_slide("Empirical anchors (scale, not targets)", "B6")
tf = _box(s, 0.9, 2.3, 11.5, 3.4)
bullets(tf, [
    [("•  ", RED), ("SVB, March 2023", INK), (": $42B withdrawn in a single day — fastest run in modern "
     "history; depositors were networked startups watching the same feeds.", BODY)],
    [("•  ", RED), ("Iyer & Puri (2012)", INK), (": in a real run, only ~3–7% of depositors withdrawing "
     "can threaten a bank — you never need a majority.", BODY)],
], size=20, space_after=14)
tf2 = _box(s, 0.9, 5.4, 11.5, 1.0)
_para(tf2, [("We use these for a sense of scale — ", INK), ("not", RED), (" as calibration targets. "
            "We're characterizing a simulation, not forecasting a real bank.", INK)], 18, first=True, line=1.3)

# B6 — Scope & limitations
s = backup_slide([("Scope & limitations ", INK), ("(what we don't claim)", MUTED)], "B7")
tf = _box(s, 0.9, 2.3, 11.5, 4.2)
bullets(tf, [
    [("•  ", RED), ("12 agents, 2 banks, simplified mechanics", INK), (" — a controlled probe, not a market model.", BODY)],
    [("•  ", RED), ("Some findings are single-pair probes", INK), (" (esp. the Central Bank) — labelled as "
     "probes, not statistics.", BODY)],
    [("•  ", RED), ("LLM calls are non-deterministic", INK), (" — we characterize patterns and re-run, rather "
     "than report a point estimate.", BODY)],
    [("•  ", RED), ("Not a forecast", INK), (": every claim is “in our simulation, these agents did X” "
     "— never “real bank runs will be X.”", BODY)],
], size=20, space_after=13)

# B7 — Central Bank detail
s = backup_slide([("Central Bank detail — ", INK), ("AI vs. rule-based", RED)], "B8")
tf = _box(s, 0.9, 2.3, 11.5, 3.6)
bullets(tf, [
    [("•  False-alarm cascade (10/12 ran on a ", BODY), ("healthy", INK), (" bank), two regulators.", BODY)],
    [("•  ", RED), ("AI regulator", INK), (": read reserves (~30%) + that zero agents had actually "
     "withdrawn yet → chose ", BODY), ("do nothing", INK), (" — correctly.", BODY)],
    [("•  ", RED), ("Rule-based", INK), (": fired a blanket deposit guarantee regardless of bank health.", BODY)],
], size=20, space_after=13)
tf2 = _box(s, 0.9, 5.5, 11.5, 1.1)
_para(tf2, [("Caveat: a single pair of runs. The robust point isn't “AI is smarter” — it's that the "
            "cascade window is ", INK), ("seconds", RED), (", so any human-paced regulator is too slow "
            "regardless.", INK)], 18, first=True, line=1.3)

# ── Speaker notes (the TIGHT script — shows in PowerPoint Presenter View) ─────
# Mirrors PRESENTATION_SCRIPT_TIGHT.md. Edit both if the talk track changes.
NOTES = {
1: (
"[0:00 · HOOK]\n"
"In five years, most of us won't move our own money in a crisis. An AI agent will — "
"checking our accounts, reading the news, deciding whether to pull our savings, in "
"seconds, while we're asleep.\n"
"So we built a sandbox to ask one question: when a bank rumor hits and the depositors "
"are AI agents instead of people — what happens?"
),
2: (
"[0:20 · PREMISE]\n"
"A human bank run is slow because noticing is slow. You have to see the news, sit with "
"it, second-guess, maybe call someone — and most people aren't watching at 3am. Even "
"SVB, which ran entirely online, took most of a day. That hesitation is part of what "
"stops a rumor from becoming a collapse.\n"
"AI delegates remove it. They watch every feed, act in seconds, and tend to reason "
"alike. We wanted to see what that does to a run."
),
3: (
"[0:35 · WHAT THIS IS]\n"
"This is not a forecast. It's twelve AI agents, each making real LLM calls to decide "
"what to do with someone's money. Each has a persona — a cautious retiree, a trader, a "
"gig worker, an institutional treasurer. They hold deposits at two banks. A rumor hits "
"Bank A. Each sees the rumor, sees what others do, and decides: hold, partially "
"withdraw, or run. Everything I show you is the actual reasoning the models produced."
),
4: (
"[1:00 · LIVE DEMO — switch to dashboard]\n\n"
"LOAD A RUN (Presets page):\n"
"\"I'll load a run we did earlier so we're not waiting on API calls — but every "
"decision here was made live.\"\n"
"[CLICK] Presets -> Load saved run -> \"Strong rumor · Bank healthy · AI speed.\"\n"
"\"Note the setup: a strong, alarming rumor about Bank A — but the bank is fine. It's "
"solvent. No one's money is actually at risk.\"  [CLICK Load and view ->]\n\n"
"THE SUMMARY (Inspect page):\n"
"\"Look at the top strip: this is what twelve agents did with a false alarm. All twelve "
"tried to exit — on a healthy bank.\"\n\n"
"READ A MIND (spend time here):\n"
"[CLICK] the first mover, or an Institutional Treasurer.\n"
"\"This is real model output, not a script. It weighs the withdrawal fee against the "
"risk to principal, decides the fee is worth it, and runs — on a bank that was fine.\"\n"
"[CLICK] a Cautious Retiree.\n"
"\"Same rumor, different persona — and listen to the voice. It reasons like a frightened "
"retiree, not a trader. You can tell the personas apart without the labels. That's the "
"heterogeneity doing real work. This is the part I'd never seen before: you can read "
"what the AI was thinking when it decided to run.\"\n\n"
"[2:50 · THREE FINDINGS — Findings page]\n"
"\"We ran dozens of these. Three patterns held up. Top to bottom.\"\n\n"
"FINDING 1 — Same personas, different brains (THE HEADLINE):\n"
"\"This one surprised us. We froze everything — same personas, same rumor, same healthy "
"bank — and changed one thing: which model decides. Withdrawing here is the wrong call; "
"the bank is fine. With Claude, about half run. With GPT, every one runs. The model "
"alone takes it from a scare to a full collapse — and these bars are replicated four to "
"five runs each, so the ranking is stable. The punchline: if everyone's money is run by "
"the same model, you don't have a diversified crowd — you have a monoculture that panics "
"in lockstep. The model itself is a systemic risk.\"\n\n"
"FINDING 2 — Speed (scroll down):\n"
"\"Second: speed. Same scenario at AI speed versus human speed — a 90-second delay per "
"agent. This is the time for half the depositors to decide, not for cash to settle: "
"about 5 seconds versus 15. Roughly three times faster. And moving money is "
"rate-limited, so the requests arrive faster than the bank can pay. It freezes before "
"most cash settles — only 5 of 12 got paid at AI speed, versus 11 of 12 at human speed. "
"Machine speed doesn't move money faster — it freezes the bank faster and locks more "
"people out.\"\n\n"
"FINDING 3 — The words decide (scroll down):\n"
"\"Third: we changed only the wording — same bank, same credibility label. Soft language "
"— 'some concern, worth monitoring' — and 31% leaves; the bank survives. Crisis language "
"— 'cannot process withdrawals' — and 93% drains in under nine seconds. The words, not "
"the credibility tag, decide whether the bank lives. (We tested more — a regulator, "
"credibility sweeps. Happy to take those in Q&A.)\""
),
5: (
"[4:20 · ANCHOR]\n"
"Is this realistic? We don't model real banks. But for scale: in March 2023, Silicon "
"Valley Bank lost $42 billion in one day — the fastest run in modern history. Why? Its "
"depositors were all startups, advised by the same VCs, watching the same threads — but "
"that took an unusual, tightly-correlated base. AI delegation manufactures that "
"correlation for any bank — and adds a new axis: the same model."
),
6: (
"[~4:50 · TWO NUMBERS]\n"
"Two numbers to remember. One: the model alone took a healthy bank from a 50% scare to a "
"100% run — nothing else changed. Two: AI delegation hit the halfway mark about 3x "
"faster than humans. Homogeneous and instant — that's the combination."
),
7: (
"[5:05 · WHY IT MATTERS]\n"
"Three things for anyone building or regulating AI money-managers:\n"
"1. The model is a systemic risk factor — a fleet on one model panics in lockstep.\n"
"2. Speed compresses the response window to seconds — human oversight can't keep up.\n"
"3. Agents react to alarming language, not stated credibility — and they're just as "
"confident when they're wrong.\n"
"This is a twelve-agent sandbox, not a forecast. The point is to find these failure "
"modes cheaply, now — before this is how finance actually works."
),
8: (
"[5:30 · SANDBOX + CLOSE]\n"
"\"We didn't hard-code these scenarios. It all runs on a sandbox where you build any "
"run. Here's a short recording.\"\n"
"[PLAY THE RECORDING — narrate lightly over the controls as they appear]\n"
"\"You describe an agent in plain English, set the bank's health and payout speed, choose "
"the population mix, write the rumor, and optionally drop in a regulator or a counter-signal "
"that talks the agents back. Name it, run it, reload it.\n"
"(As the result lands) Same engine, different scenario. And this is really "
"infrastructure: a regulator can stress-test its delegation rules before approving them, "
"a fintech can test its agent before launch — and as you feed in real data on who "
"delegates to which model, you can size the systemic risk we found here. Eventually you "
"just type the crisis and the sandbox builds it. Twelve agents is the start — this is "
"how anyone probes the next one. Thank you — happy to dig into any of it.\"\n"
"[Recording: ~30-40s, end on the Inspect view. Pre-rendered — do not run live.]"
),
}
slides_list = list(prs.slides)
for num, note in NOTES.items():
    slides_list[num - 1].notes_slide.notes_text_frame.text = note

# ── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).parent.parent / "SLIDES.pptx"
prs.save(out)
print(f"Wrote {out} · {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
